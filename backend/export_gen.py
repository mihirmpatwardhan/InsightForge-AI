import re
from io import BytesIO
from docx import Document
from docx.shared import RGBColor, Pt
from reportlab.lib.pagesizes import A4
from reportlab.lib import colors
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, HRFlowable
from pptx import Presentation
from pptx.util import Inches, Pt as PPTpt
from pptx.dml.color import RGBColor as PPTRGBColor
from pptx.enum.text import PP_ALIGN


def md_to_html(text: str) -> str:
    text = re.sub(r"\*\*(.*?)\*\*", r"<b>\1</b>", text)
    return text


def create_docx_bytes(title: str, content: str) -> bytes:
    buffer = BytesIO()
    doc = Document()
    h = doc.add_heading(level=0)
    run = h.add_run(title)
    run.font.color.rgb = RGBColor(79, 70, 229)
    run.font.size = Pt(20)

    for line in content.split("\n"):
        line = line.strip()
        if not line:
            continue
        if line.startswith("# "):
            doc.add_heading(line[2:].replace("**", ""), level=1)
        elif line.startswith("## "):
            doc.add_heading(line[3:].replace("**", ""), level=2)
        elif line.startswith("### "):
            doc.add_heading(line[4:].replace("**", ""), level=3)
        elif line.startswith(("- ", "* ")):
            p = doc.add_paragraph(style="List Bullet")
            _add_bold_runs(p, line[2:])
        elif re.match(r"^\d+\.\s", line):
            p = doc.add_paragraph(style="List Number")
            _add_bold_runs(p, re.sub(r"^\d+\.\s", "", line))
        else:
            p = doc.add_paragraph()
            _add_bold_runs(p, line)

    doc.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()


def _add_bold_runs(p, text: str):
    parts = re.split(r"(\*\*.*?\*\*)", text)
    for part in parts:
        if part.startswith("**") and part.endswith("**"):
            run = p.add_run(part[2:-2])
            run.bold = True
        else:
            p.add_run(part)


def create_pdf_bytes(title: str, content: str) -> bytes:
    buffer = BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=A4, rightMargin=40, leftMargin=40, topMargin=40, bottomMargin=40)
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle("DocTitle", parent=styles["Heading1"], fontName="Helvetica-Bold", fontSize=20, textColor=colors.HexColor("#4f46e5"), spaceAfter=12)
    h1_style = ParagraphStyle("H1", parent=styles["Heading2"], fontName="Helvetica-Bold", fontSize=13, textColor=colors.HexColor("#1e1b4b"), spaceBefore=10, spaceAfter=4)
    body_style = ParagraphStyle("Body", parent=styles["Normal"], fontName="Helvetica", fontSize=9.5, leading=14, textColor=colors.HexColor("#334155"), spaceAfter=5)
    bullet_style = ParagraphStyle("Bullet", parent=body_style, leftIndent=14, spaceAfter=3)

    story = [
        Paragraph(md_to_html(title), title_style),
        HRFlowable(width="100%", thickness=1, color=colors.HexColor("#6366f1"), spaceAfter=12)
    ]

    for line in content.split("\n"):
        line = line.strip()
        if not line:
            story.append(Spacer(1, 4))
            continue
        html = md_to_html(line)
        if line.startswith(("# ", "## ", "### ")):
            story.append(Paragraph(re.sub(r"^#+\s*", "", html), h1_style))
        elif line.startswith(("- ", "* ")):
            story.append(Paragraph(f"• {re.sub(r'^[-*]\\s*', '', html)}", bullet_style))
        elif re.match(r"^\d+\.\s", line):
            story.append(Paragraph(html, bullet_style))
        else:
            story.append(Paragraph(html, body_style))

    doc.build(story)
    buffer.seek(0)
    return buffer.getvalue()


def create_pptx_bytes(title: str, content: str, dataset_name: str = "Dataset") -> bytes:
    """
    Creates a branded PowerPoint presentation from markdown-formatted AI content.
    """
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    BG_COLOR = PPTRGBColor(6, 9, 17)         # #060911
    ACCENT = PPTRGBColor(99, 102, 241)        # #6366f1
    TEXT_LIGHT = PPTRGBColor(225, 232, 240)   # #e1e8f0
    TEXT_MUTED = PPTRGBColor(148, 163, 184)   # #94a3b8
    HIGHLIGHT = PPTRGBColor(165, 180, 252)    # #a5b4fc

    def set_slide_bg(slide, color: PPTRGBColor):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_text_box(slide, text, left, top, width, height,
                     font_size=18, bold=False, color=TEXT_LIGHT, align=PP_ALIGN.LEFT):
        txb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = PPTpt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color

    def add_accent_line(slide, top):
        from pptx.util import Emu
        line = slide.shapes.add_shape(1, Inches(0.5), Inches(top), Inches(12.3), Emu(36000))
        line.fill.solid()
        line.fill.fore_color.rgb = ACCENT
        line.line.fill.background()

    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide1, BG_COLOR)
    add_accent_line(slide1, 2.0)
    add_text_box(slide1, "NexusViz AI", 0.5, 0.4, 12, 0.8, font_size=14, color=HIGHLIGHT, align=PP_ALIGN.LEFT)
    add_text_box(slide1, title, 0.5, 2.2, 12, 1.2, font_size=36, bold=True, color=TEXT_LIGHT, align=PP_ALIGN.LEFT)
    add_text_box(slide1, f"Dataset: {dataset_name}", 0.5, 3.6, 12, 0.6, font_size=16, color=TEXT_MUTED)
    add_text_box(slide1, "AI-Generated Executive Report  •  Confidential", 0.5, 6.8, 12, 0.5, font_size=11, color=TEXT_MUTED, align=PP_ALIGN.LEFT)

    # Parse content into sections
    sections = []
    current_section = None
    current_bullets = []

    for line in content.split("\n"):
        line = line.strip()
        if not line:
            continue
        if line.startswith(("# ", "## ", "### ")):
            if current_section:
                sections.append((current_section, current_bullets))
            current_section = re.sub(r"^#+\s*", "", line).replace("**", "")
            current_bullets = []
        elif line.startswith(("- ", "* ")):
            current_bullets.append(("bullet", re.sub(r"^[-*]\s*", "", line).replace("**", "")))
        elif re.match(r"^\d+\.\s", line):
            current_bullets.append(("numbered", re.sub(r"^\d+\.\s", "", line).replace("**", "")))
        else:
            current_bullets.append(("text", line.replace("**", "")))

    if current_section:
        sections.append((current_section, current_bullets))

    # Generate a slide per section (max 10 slides)
    for i, (section_title, bullets) in enumerate(sections[:10]):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(slide, BG_COLOR)
        add_accent_line(slide, 1.5)

        # Slide number badge
        add_text_box(slide, f"{i+2:02d}", 12.3, 0.2, 0.8, 0.5, font_size=11, color=ACCENT)

        # Section title
        add_text_box(slide, section_title, 0.5, 0.3, 11.5, 1.0, font_size=26, bold=True, color=HIGHLIGHT)

        # Bullets (max 8 per slide)
        y_pos = 1.8
        for kind, text in bullets[:8]:
            prefix = "•  " if kind == "bullet" else ("→  " if kind == "numbered" else "")
            safe_text = (prefix + text)[:200]
            add_text_box(slide, safe_text, 0.7, y_pos, 11.5, 0.55, font_size=15, color=TEXT_LIGHT)
            y_pos += 0.6

    # Last slide: Thank You
    slide_end = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide_end, BG_COLOR)
    add_accent_line(slide_end, 3.2)
    add_text_box(slide_end, "Thank You", 0.5, 2.8, 12, 1.2, font_size=42, bold=True, color=HIGHLIGHT, align=PP_ALIGN.CENTER)
    add_text_box(slide_end, "Generated by NexusViz AI  •  Enterprise Data Intelligence Platform", 0.5, 4.2, 12, 0.6, font_size=13, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()
